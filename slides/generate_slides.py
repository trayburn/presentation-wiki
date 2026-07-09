import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors
    bg_color = RGBColor(15, 23, 42)      # Slate 900
    text_light = RGBColor(248, 250, 252) # Slate 50
    text_muted = RGBColor(148, 163, 184) # Slate 400
    accent_sky = RGBColor(56, 189, 248)  # Sky 400
    accent_indigo = RGBColor(129, 140, 248) # Indigo 400

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_title(slide, text, subtitle_text=None):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = text_light
        p.font.name = "Arial"
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(18)
            p2.font.color.rgb = text_muted
            p2.font.name = "Arial"
            p2.space_before = Pt(8)

    # ----------------------------------------------------
    # Slide 1: Title
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "SURVIVING THE AI DELUGE"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = accent_sky
    p.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    p2.text = "Architecting an Agentic Memory with OKF and Obsidian"
    p2.font.size = Pt(26)
    p2.font.color.rgb = text_light
    p2.font.name = "Arial"
    p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = "Tim Rayburn | Vice President, Improving"
    p3.font.size = Pt(20)
    p3.font.color.rgb = text_muted
    p3.font.name = "Arial"
    p3.space_before = Pt(32)

    # ----------------------------------------------------
    # Slide 2: The Problem: The AI Information Deluge
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_title(slide2, "The Problem: The AI Information Deluge", "Why trying to manually process the flood of AI changes fails developers")

    # Left Column: The Content Flood
    col1_box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True
    
    p_col1_header = tf1.paragraphs[0]
    p_col1_header.text = "1. The Information Deluge"
    p_col1_header.font.size = Pt(22)
    p_col1_header.font.bold = True
    p_col1_header.font.color.rgb = accent_indigo
    p_col1_header.font.name = "Arial"
    p_col1_header.space_after = Pt(12)

    bullet1 = tf1.add_paragraph()
    bullet1.text = "• A weekly deluge of new models, benchmark papers, and API updates."
    bullet1.font.size = Pt(16)
    bullet1.font.color.rgb = text_light
    bullet1.font.name = "Arial"
    bullet1.space_after = Pt(10)

    bullet2 = tf1.add_paragraph()
    bullet2.text = "• Bookmarks and traditional wikis fail because manual curation takes too long."
    bullet2.font.size = Pt(16)
    bullet2.font.color.rgb = text_light
    bullet2.font.name = "Arial"
    bullet2.space_after = Pt(10)

    bullet3 = tf1.add_paragraph()
    bullet3.text = "• Tuning out is risky—ignoring token economics or model limits costs real money."
    bullet3.font.size = Pt(16)
    bullet3.font.color.rgb = text_light
    bullet3.font.name = "Arial"

    # Right Column: The Cognitive Tax
    col2_box = slide2.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.6), Inches(4.5))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True
    
    p_col2_header = tf2.paragraphs[0]
    p_col2_header.text = "2. The Cognitive Tax"
    p_col2_header.font.size = Pt(22)
    p_col2_header.font.bold = True
    p_col2_header.font.color.rgb = accent_indigo
    p_col2_header.font.name = "Arial"
    p_col2_header.space_after = Pt(12)

    bullet_r1 = tf2.add_paragraph()
    bullet_r1.text = "• Cognitive Surrender: Overwhelmed users stop reasoning and accept faulty outputs."
    bullet_r1.font.size = Pt(16)
    bullet_r1.font.color.rgb = text_light
    bullet_r1.font.name = "Arial"
    bullet_r1.space_after = Pt(10)

    bullet_r2 = tf2.add_paragraph()
    bullet_r2.text = "• Cognitive Debt: Relying on AI generation without comprehension degrades memory."
    bullet_r2.font.size = Pt(16)
    bullet_r2.font.color.rgb = text_light
    bullet_r2.font.name = "Arial"
    bullet_r2.space_after = Pt(10)

    bullet_r3 = tf2.add_paragraph()
    bullet_r3.text = "• Need: A structured filter to separate high-value signals from marketing noise."
    bullet_r3.font.size = Pt(16)
    bullet_r3.font.color.rgb = text_light
    bullet_r3.font.name = "Arial"

    # ----------------------------------------------------
    # Slide 3: The Pattern: Living Context
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_title(slide3, "The Pattern: Living Context", "Using agents to filter, summarize, and structure the flood of information")

    box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True

    p_intro = tf.paragraphs[0]
    p_intro.text = "Instead of manual curation, deploy AI agents to manage your team's knowledge flow."
    p_intro.font.size = Pt(20)
    p_intro.font.bold = True
    p_intro.font.color.rgb = accent_sky
    p_intro.font.name = "Arial"
    p_intro.space_after = Pt(24)

    steps = [
        ("1. Automated Ingestion", "AI agents parse incoming raw PDFs, blog posts, and study clippings on arrival, extracting key developer impact points."),
        ("2. Open Knowledge Format (OKF)", "Write insights into clean, version-controlled markdown notes with metadata headers. No custom database needed."),
        ("3. Obsidian Graph Curation", "Use Obsidian's auto-generated semantic graph to map dependencies, review linkages, and audit the ingested memory.")
    ]

    for title, desc in steps:
        p_step = tf.add_paragraph()
        p_step.text = f"{title}: {desc}"
        p_step.font.size = Pt(16)
        p_step.font.color.rgb = text_light
        p_step.font.name = "Arial"
        p_step.space_after = Pt(16)

    # ----------------------------------------------------
    # Slide 4: Open Knowledge Format (OKF)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_title(slide4, "Open Knowledge Format (OKF)", "A minimal, human-readable, and agent-friendly knowledge spec")

    col1_box = slide4.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True

    p_c1 = tf1.paragraphs[0]
    p_c1.text = "OKF Architecture"
    p_c1.font.size = Pt(22)
    p_c1.font.bold = True
    p_c1.font.color.rgb = accent_sky
    p_c1.font.name = "Arial"
    p_c1.space_after = Pt(14)

    items = [
        ("Directory of Markdown", "No proprietary database or schema registries needed."),
        ("YAML Frontmatter", "Forces metadata discipline: 'type', 'title', 'tags', 'timestamp'."),
        ("Citations & Links", "Standard Markdown relative links construct the semantic graph.")
    ]

    for item_t, item_d in items:
        p_item = tf1.add_paragraph()
        p_item.text = f"• {item_t}: {item_d}"
        p_item.font.size = Pt(15)
        p_item.font.color.rgb = text_light
        p_item.font.name = "Arial"
        p_item.space_after = Pt(12)

    col2_box = slide4.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.5))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True

    p_c2 = tf2.paragraphs[0]
    p_c2.text = "Sample Concept Document"
    p_c2.font.size = Pt(18)
    p_c2.font.bold = True
    p_c2.font.color.rgb = text_muted
    p_c2.font.name = "Arial"
    p_c2.space_after = Pt(10)

    p_code = tf2.add_paragraph()
    p_code.text = (
        "---\n"
        "type: Reference\n"
        "title: BCG & Harvard AI Field Study\n"
        "tags: [ai-adoption, research]\n"
        "timestamp: 2026-06-29T11:50:00Z\n"
        "---\n"
        "# Summary\n"
        "Study of 758 consultants showing GPT-4 quality boost (40%)\n"
        "vs. out-of-boundary quality drop (19 percentage points).\n\n"
        "See also [MIT Study](/sources/mit-cognitive-debt.md)."
    )
    p_code.font.size = Pt(13)
    p_code.font.color.rgb = text_light
    p_code.font.name = "Courier New"

    # ----------------------------------------------------
    # Slide 5: The Ingestion Canon (Case Studies)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_title(slide5, "The Ingestion Canon: The Data We Cured", "Key studies loaded into our demonstration knowledge archive")

    box = slide5.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True

    canon = [
        ("1. BCG & Harvard Study (SSRN 2023)", "Analyzed 758 consultants. AI boosted in-frontier quality by 40%, but decreased out-of-frontier accuracy by 19 percentage points when trust was uncalibrated."),
        ("2. MIT Media Lab Study (arXiv 2025)", "Measured cognitive debt. EEG scans showed reduced brain activity during AI-assisted drafting, and 83% of ChatGPT writers could not recall what they wrote."),
        ("3. Wharton Study on Cognitive Surrender (2025)", "Analyzed 1,372 participants. High trust led users to stop constructing answers entirely—yielding an 80% follow-rate on faulty AI outputs."),
        ("4. Devlin Liles' 'The Meter Picks the Winners' (2026)", "Explored tokenizer repricing, high per-user consumption variance, and building local owned floors (DGX Sparks) to anchor model routing.")
    ]

    is_first = True
    for c_title, c_desc in canon:
        p_c = tf.paragraphs[0] if is_first else tf.add_paragraph()
        p_c.text = f"{c_title}"
        p_c.font.size = Pt(16)
        p_c.font.bold = True
        p_c.font.color.rgb = accent_indigo
        p_c.font.name = "Arial"
        p_c.space_after = Pt(2)
        if not is_first:
            p_c.space_before = Pt(12)
        is_first = False
        
        p_d = tf.add_paragraph()
        p_d.text = f"   {c_desc}"
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = text_light
        p_d.font.name = "Arial"

    # ----------------------------------------------------
    # Slide 6: LIVE DEMO: Taming the Deluge
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_title(slide6, "LIVE DEMO: Taming the Deluge", "Let's look at the tools...")

    box = slide6.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True

    p_demo_title = tf.paragraphs[0]
    p_demo_title.text = "Interactive Walkthrough Outline:"
    p_demo_title.font.size = Pt(22)
    p_demo_title.font.bold = True
    p_demo_title.font.color.rgb = accent_sky
    p_demo_title.font.name = "Arial"
    p_demo_title.space_after = Pt(16)

    demo_steps = [
        "1. Ingestion: Run a Python script that takes a raw PDF clipping and writes a structured OKF document.",
        "2. Visualization: Open the resulting folder directly in Obsidian to view the frontmatter and updated graph view.",
        "3. Cross-Linking: Watch how Obsidian's internal linking mirrors the agent's semantic relationships.",
        "4. Consumption: Execute a local RAG agent that traverses OKF links to answer architectural and cost queries."
    ]

    for step in demo_steps:
        p_step = tf.add_paragraph()
        p_step.text = step
        p_step.font.size = Pt(16)
        p_step.font.color.rgb = text_light
        p_step.font.name = "Arial"
        p_step.space_after = Pt(12)

    # ----------------------------------------------------
    # Slide 7: Monday Morning Roadmap
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_title(slide7, "Monday Morning Roadmap", "How to build your team's AI information filter")

    box = slide7.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True

    roadmap_steps = [
        ("Initialize a Shared Git Repo", "Establish a minimal OKF folder structure for your team's architecture, specs, and reference materials. Treat it as code."),
        ("Automate Ingestion Gardens", "Set up simple post-commit hooks or scheduler tasks that convert Slack decisions, transcripts, and meeting summaries into OKF."),
        ("Enforce Routing Rubrics", "Use your team's knowledge index as a routing rulebook—direct routine questions to cheap local/open models, saving premium APIs for frontier reasoning."),
        ("Defend Against Cognitive Debt", "Maintain a 'nemesis review' posture. Audit AI-generated code and specs with higher skepticism than human colleague contributions.")
    ]

    is_first = True
    for r_title, r_desc in roadmap_steps:
        p_r = tf.paragraphs[0] if is_first else tf.add_paragraph()
        p_r.text = f"• {r_title}"
        p_r.font.size = Pt(18)
        p_r.font.bold = True
        p_r.font.color.rgb = accent_sky
        p_r.font.name = "Arial"
        p_r.space_after = Pt(2)
        if not is_first:
            p_r.space_before = Pt(14)
        is_first = False

        p_rd = tf.add_paragraph()
        p_rd.text = f"  {r_desc}"
        p_rd.font.size = Pt(15)
        p_rd.font.color.rgb = text_light
        p_rd.font.name = "Arial"

    # Save
    out_dir = "/Users/bob/vaults/Improving/01-Projects/Public Speaking/03 - Complete & Submittable/Living Context - Architecting Agentic Memory"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "Living Context.pptx")
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")

if __name__ == "__main__":
    create_presentation()

